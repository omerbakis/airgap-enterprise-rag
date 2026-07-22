"""PDF/DOCX/TXT/MD dosyalarını ortak bir Block listesine ayrıştırır.

Her parser, dokümanın yapısını (başlık seviyesi, paragraf, tablo, liste öğesi)
olabildiğince koruyarak sıralı bir Block listesi üretir; chunker (bkz.
ingestion/chunker.py) bu listeyi section-aware chunk'lara dönüştürür.
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

import fitz  # pymupdf
import openpyxl
from bs4 import BeautifulSoup, Tag
from docx import Document as DocxDocument
from docx.oxml.ns import qn
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph
from pptx import Presentation

BlockType = str  # "heading" | "paragraph" | "table" | "list_item"


@dataclass
class Block:
    type: BlockType
    text: str
    level: int | None = None  # yalnızca heading için (1=H1, 2=H2, ...)
    page: int | None = None


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".md", ".xlsx", ".pptx", ".html", ".htm"}


def parse_file(path: Path) -> list[Block]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _parse_pdf(path)
    if ext == ".docx":
        return _parse_docx(path)
    if ext in (".txt", ".md"):
        return _parse_text(path)
    if ext == ".xlsx":
        return _parse_xlsx(path)
    if ext == ".pptx":
        return _parse_pptx(path)
    if ext in (".html", ".htm"):
        return _parse_html(path)
    raise ValueError(f"Desteklenmeyen dosya türü: {ext}")


# --------------------------------------------------------------------------
# TXT / Markdown
# --------------------------------------------------------------------------

_MD_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)")
_MD_LIST_RE = re.compile(r"^\s*([-*+]|\d+[.)])\s+")
_MD_TABLE_ROW_RE = re.compile(r"^\s*\|.*\|\s*$")
_MD_TABLE_SEP_RE = re.compile(r"^[\s|:-]+$")


def _is_markdown_table(lines: list[str]) -> bool:
    if len(lines) < 2:
        return False
    return bool(
        _MD_TABLE_ROW_RE.match(lines[0])
        and _MD_TABLE_SEP_RE.match(lines[1])
        and "-" in lines[1]
    )


def _parse_text(path: Path) -> list[Block]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    blocks: list[Block] = []
    for para in re.split(r"\n\s*\n", raw):
        para = para.strip()
        if not para:
            continue
        heading_match = _MD_HEADING_RE.match(para)
        if heading_match:
            blocks.append(Block(type="heading", level=len(heading_match.group(1)), text=heading_match.group(2).strip()))
            continue
        lines = [line for line in para.splitlines() if line.strip()]
        if _is_markdown_table(lines):
            blocks.append(Block(type="table", text="\n".join(lines)))
            continue
        if lines and all(_MD_LIST_RE.match(line) for line in lines):
            for line in lines:
                item_text = _MD_LIST_RE.sub("", line).strip()
                blocks.append(Block(type="list_item", text=item_text))
            continue
        blocks.append(Block(type="paragraph", text=para))
    return blocks


# --------------------------------------------------------------------------
# DOCX
# --------------------------------------------------------------------------

_HEADING_STYLE_RE = re.compile(r"^Heading (\d+)$", re.IGNORECASE)


def _iter_docx_block_items(document: DocxDocument):
    parent_elm = document.element.body
    for child in parent_elm.iterchildren():
        if child.tag == qn("w:p"):
            yield DocxParagraph(child, document)
        elif child.tag == qn("w:tbl"):
            yield DocxTable(child, document)


def _docx_table_to_markdown(table: DocxTable) -> str:
    rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    if not rows:
        return ""
    header, *body_rows = rows
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * len(header)) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in body_rows)
    return "\n".join(lines)


def _parse_docx(path: Path) -> list[Block]:
    document = DocxDocument(str(path))
    blocks: list[Block] = []
    for item in _iter_docx_block_items(document):
        if isinstance(item, DocxTable):
            md = _docx_table_to_markdown(item)
            if md:
                blocks.append(Block(type="table", text=md))
            continue

        # python-docx, bir paragraf içindeki satır sonlarını (<w:br/>) "\n" olarak
        # metne dahil eder; bu, tek satırlık olması beklenen başlık/section_path
        # metadata'sına çok satırlı içerik sızmasına yol açar — bu yüzden tüm
        # iç boşluklar tek boşluğa normalize edilir.
        text = " ".join(item.text.split())
        if not text:
            continue

        style_name = item.style.name if item.style is not None else ""
        heading_match = _HEADING_STYLE_RE.match(style_name or "")
        if heading_match:
            blocks.append(Block(type="heading", level=int(heading_match.group(1)), text=text))
            continue
        if style_name and "List" in style_name:
            blocks.append(Block(type="list_item", text=text))
            continue
        blocks.append(Block(type="paragraph", text=text))
    return blocks


# --------------------------------------------------------------------------
# PDF
# --------------------------------------------------------------------------
# Not: PDF'lerde gerçek stil bilgisi yok; başlıklar, sayfadaki en sık görülen
# ("gövde") font boyutuna göre bağıl olarak büyük satırlar üzerinden
# sezgisel (heuristic) olarak tespit edilir. Tablo tespiti bu sürümde
# kapsam dışıdır.


def _parse_pdf(path: Path) -> list[Block]:
    doc = fitz.open(str(path))
    line_sizes: list[float] = []
    pages_lines: list[list[tuple[str, float, int]]] = []

    for page_index, page in enumerate(doc, start=1):
        page_lines: list[tuple[str, float, int]] = []
        page_dict = page.get_text("dict")
        for block in page_dict.get("blocks", []):
            for line in block.get("lines", []):
                spans = line.get("spans", [])
                if not spans:
                    continue
                text = "".join(span.get("text", "") for span in spans).strip()
                if not text:
                    continue
                size = max(span.get("size", 0.0) for span in spans)
                page_lines.append((text, size, page_index))
                line_sizes.append(size)
        pages_lines.append(page_lines)
    doc.close()

    if not line_sizes:
        return []

    body_size = _most_common_size(line_sizes)
    blocks: list[Block] = []
    paragraph_buffer: list[str] = []
    buffer_page: int | None = None

    def flush_paragraph():
        if paragraph_buffer:
            blocks.append(Block(type="paragraph", text=" ".join(paragraph_buffer), page=buffer_page))
            paragraph_buffer.clear()

    for page_lines in pages_lines:
        for text, size, page_index in page_lines:
            level = _heading_level_for_size(size, body_size)
            if level is not None:
                flush_paragraph()
                blocks.append(Block(type="heading", level=level, text=text, page=page_index))
                continue
            if not paragraph_buffer:
                buffer_page = page_index
            paragraph_buffer.append(text)
        flush_paragraph()

    return blocks


def _most_common_size(sizes: list[float]) -> float:
    rounded = [round(s, 1) for s in sizes]
    return max(set(rounded), key=rounded.count)


def _heading_level_for_size(size: float, body_size: float) -> int | None:
    if body_size <= 0:
        return None
    ratio = size / body_size
    if ratio >= 1.4:
        return 1
    if ratio >= 1.15:
        return 2
    return None


# --------------------------------------------------------------------------
# XLSX
# --------------------------------------------------------------------------
# Her sayfa bir bölüm (heading) olur; satırlar, tek bir dev/sınırsız chunk
# oluşmasını önlemek için sabit boyutlu gruplar halinde ayrı tablo Block'larına
# bölünür (DOCX/MD tabloları küçük olduğundan atomic kalır, ama tablo kaynaklı
# dev chunk riski en çok elektronik tablolarda).

_XLSX_ROWS_PER_CHUNK = 30


def _rows_to_markdown_table(header: list[str], rows: list[list[str]]) -> str:
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * len(header)) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in rows)
    return "\n".join(lines)


def _parse_xlsx(path: Path) -> list[Block]:
    workbook = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    blocks: list[Block] = []
    try:
        for sheet in workbook.worksheets:
            rows = [
                ["" if cell is None else str(cell) for cell in row]
                for row in sheet.iter_rows(values_only=True)
            ]
            rows = [row for row in rows if any(cell.strip() for cell in row)]
            if not rows:
                continue
            blocks.append(Block(type="heading", level=2, text=sheet.title))
            header, *data_rows = rows
            if not data_rows:
                continue
            for start in range(0, len(data_rows), _XLSX_ROWS_PER_CHUNK):
                chunk_rows = data_rows[start : start + _XLSX_ROWS_PER_CHUNK]
                blocks.append(Block(type="table", text=_rows_to_markdown_table(header, chunk_rows)))
    finally:
        workbook.close()
    return blocks


# --------------------------------------------------------------------------
# PPTX
# --------------------------------------------------------------------------
# Slaytlar genelde madde işaretli kısa metinlerden oluştuğu için başlık dışı
# tüm metin kutuları list_item olarak ele alınır (prosedür/madde tarzı
# atomic chunklama ile tutarlı, bkz. chunker.py).


def _parse_pptx(path: Path) -> list[Block]:
    presentation = Presentation(str(path))
    blocks: list[Block] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title_text = title_shape.text.strip() if title_shape is not None and title_shape.has_text_frame else ""
        blocks.append(Block(type="heading", level=2, text=title_text or f"Slayt {slide_index}"))

        for shape in slide.shapes:
            if shape is title_shape or not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    blocks.append(Block(type="list_item", text=text))
    return blocks


# --------------------------------------------------------------------------
# HTML
# --------------------------------------------------------------------------
# <script>/<style> içeriği tamamen atılır (hem gürültü hem prompt-injection
# yüzeyi). Tablo içindeki hücreler ayrıca paragraf/liste olarak tekrar
# işlenmez (find_parent("table") kontrolü).


def _html_table_to_markdown(table: Tag) -> str:
    rows: list[list[str]] = []
    for tr in table.find_all("tr"):
        cells = [" ".join(cell.get_text(" ", strip=True).split()) for cell in tr.find_all(["td", "th"])]
        if cells:
            rows.append(cells)
    if not rows:
        return ""
    header, *body_rows = rows
    return _rows_to_markdown_table(header, body_rows)


def _parse_html(path: Path) -> list[Block]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()

    blocks: list[Block] = []
    for element in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "table", "li"]):
        if element.name == "table":
            md = _html_table_to_markdown(element)
            if md:
                blocks.append(Block(type="table", text=md))
            continue
        if element.find_parent("table") is not None:
            continue

        text = " ".join(element.get_text(" ", strip=True).split())
        if not text:
            continue
        if element.name == "li":
            blocks.append(Block(type="list_item", text=text))
        elif element.name == "p":
            blocks.append(Block(type="paragraph", text=text))
        else:
            blocks.append(Block(type="heading", level=int(element.name[1]), text=text))
    return blocks
