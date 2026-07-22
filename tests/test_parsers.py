from pathlib import Path

import openpyxl
from pptx import Presentation

from local_rag.ingestion.parsers import _XLSX_ROWS_PER_CHUNK, parse_file


def test_markdown_table_is_tagged_as_table_block(tmp_path):
    content = (
        "# Başlık\n\n"
        "Bir paragraf.\n\n"
        "| Kıdem | İzin |\n"
        "| --- | --- |\n"
        "| 1-5 | 14 |\n\n"
        "Tablodan sonraki paragraf.\n"
    )
    path = tmp_path / "doc.md"
    path.write_text(content, encoding="utf-8")

    blocks = parse_file(path)
    types = [b.type for b in blocks]
    assert types == ["heading", "paragraph", "table", "paragraph"]
    table_block = next(b for b in blocks if b.type == "table")
    assert table_block.text.startswith("| Kıdem | İzin |")


def test_markdown_list_is_not_misdetected_as_table(tmp_path):
    content = "- Adım 1\n- Adım 2\n"
    path = tmp_path / "doc.md"
    path.write_text(content, encoding="utf-8")

    blocks = parse_file(path)
    assert all(b.type == "list_item" for b in blocks)


def test_xlsx_sheet_becomes_heading_plus_table_blocks(tmp_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "İzin Süreleri"
    ws.append(["Kıdem", "Gün"])
    ws.append(["1-5", 14])
    ws.append(["5-15", 20])
    path = tmp_path / "tablo.xlsx"
    wb.save(str(path))

    blocks = parse_file(path)
    assert blocks[0].type == "heading"
    assert blocks[0].text == "İzin Süreleri"
    table_blocks = [b for b in blocks if b.type == "table"]
    assert len(table_blocks) == 1
    assert "Kıdem" in table_blocks[0].text
    assert "1-5" in table_blocks[0].text and "14" in table_blocks[0].text


def test_xlsx_large_sheet_splits_into_multiple_table_chunks(tmp_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Büyük Tablo"
    ws.append(["Kolon"])
    row_count = _XLSX_ROWS_PER_CHUNK + 5
    for i in range(row_count):
        ws.append([f"deger-{i}"])
    path = tmp_path / "buyuk.xlsx"
    wb.save(str(path))

    blocks = parse_file(path)
    table_blocks = [b for b in blocks if b.type == "table"]
    assert len(table_blocks) == 2
    assert "deger-0" in table_blocks[0].text
    assert f"deger-{row_count - 1}" in table_blocks[-1].text


def test_pptx_slide_becomes_heading_plus_list_items(tmp_path):
    presentation = Presentation()
    layout = presentation.slide_layouts[1]  # Title and Content
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = "Güvenlik Politikası"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "VPN zorunludur"
    p2 = tf.add_paragraph()
    p2.text = "MFA zorunludur"
    path = tmp_path / "sunum.pptx"
    presentation.save(str(path))

    blocks = parse_file(path)
    assert blocks[0].type == "heading"
    assert blocks[0].text == "Güvenlik Politikası"
    list_texts = [b.text for b in blocks if b.type == "list_item"]
    assert "VPN zorunludur" in list_texts
    assert "MFA zorunludur" in list_texts


def test_html_headings_paragraphs_lists_and_tables(tmp_path):
    content = """
    <html><body>
    <h1>Ana Başlık</h1>
    <p>Bir paragraf.</p>
    <ul><li>Madde 1</li><li>Madde 2</li></ul>
    <table>
        <tr><th>Kolon A</th><th>Kolon B</th></tr>
        <tr><td>1</td><td>2</td></tr>
    </table>
    <script>alert('tehlikeli kod çalıştırılmamalı');</script>
    </body></html>
    """
    path = tmp_path / "sayfa.html"
    path.write_text(content, encoding="utf-8")

    blocks = parse_file(path)
    types = [b.type for b in blocks]
    assert types == ["heading", "paragraph", "list_item", "list_item", "table"]
    assert blocks[0].text == "Ana Başlık"
    table_block = next(b for b in blocks if b.type == "table")
    assert "Kolon A" in table_block.text and "1" in table_block.text
    assert all("alert" not in b.text for b in blocks)
